# -*- coding: utf-8 -*-
"""가로형(Landscape 16:9) 홍보 포스터 생성 — PNG 썸네일 + 편집용 PPTX.

익명화 원칙: 개인·소속 특정 정보 일절 미포함.
실행: python docs/poster/make_poster.py  (저장소 루트 기준 어느 위치에서든 동작)
산출: docs/poster/poster.png (1920×1080) · poster.pptx · qr_app.png
"""
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib.patches import FancyBboxPatch
import qrcode

OUT = Path(__file__).resolve().parent
URL = "https://pf-aftercare-dashboard-fhdynfqdime54cstappppxqy.streamlit.app/"

# ── 팔레트 (대시보드와 동일 계열, 다크 모드 스텝) ──
BG = "#101725"; CARD = "#1A2436"
INK = "#FFFFFF"; INK2 = "#C3C2B7"; MUTED = "#8B93A3"
BLUE = "#3987E5"; ORANGE = "#EB6834"; AQUA = "#1BAF7A"; YELLOW = "#EDA100"
KR = "Apple SD Gothic Neo"
plt.rcParams["font.family"] = KR
plt.rcParams["axes.unicode_minus"] = False

# ── QR ──
qr = qrcode.QRCode(border=1, box_size=10)
qr.add_data(URL)
qr.make(fit=True)
qimg = qr.make_image(fill_color="#101725", back_color="white")
qimg.save(OUT / "qr_app.png")

# ── 캔버스 ──
fig = plt.figure(figsize=(19.2, 10.8), dpi=100)
fig.patch.set_facecolor(BG)
ax = fig.add_axes([0, 0, 1, 1]); ax.set_xlim(0, 1); ax.set_ylim(0, 1)
ax.axis("off"); ax.set_facecolor(BG)

def card(x, y, w, h, r=0.015):
    ax.add_patch(FancyBboxPatch((x, y), w, h, boxstyle=f"round,pad=0,rounding_size={r}",
                                fc=CARD, ec="none", zorder=1))

# ── 헤더 ──
ax.text(0.055, 0.945, "공개 데이터 기반 · AI 협업 프로젝트", color=AQUA,
        fontsize=17, weight="bold")
ax.text(0.052, 0.845, "부실 PF 사업장, 가격은 어떻게 움직이는가",
        color=INK, fontsize=46, weight="heavy")
ax.text(0.055, 0.785, "공개 경공매 사업장 798곳의 감정가·최저입찰가 경로를 시계열로 추적하고, "
        "점검이 필요한 사업장을 스크리닝하는 대시보드",
        color=INK2, fontsize=18.5)

# ── 스탯 스트립 ──
stats = [("798", "추적 사업장", BLUE), ("17개월", "가격 경로 관측", BLUE),
         ("129 → 65", "검증으로 확정한 변수", ORANGE), ("7,977건", "통계 검정(FDR 보정)", ORANGE)]
for i, (num, lab, c) in enumerate(stats):
    x = 0.055 + i * 0.155
    ax.text(x, 0.665, num, color=c, fontsize=34, weight="heavy")
    ax.text(x, 0.628, lab, color=MUTED, fontsize=14)

# ── 인사이트 카드 3장 ──
cards = [
    ("가격은 버티고,\n최저입찰가만 깎인다", ORANGE,
     "감정평가액은 공개기간 중 83.6% 무변동.\n공매는 감정가의 120%에서 출발해\n회차마다 -9.5%씩 인하 — 8회차엔 62%."),
    ("산업시설은 가격이 아니라\n수요의 문제", AQUA,
     "물류·공장은 시장보다 싸게 내놓아도\n체류 최장(중앙 11개월).\n감정 괴리 최소 · 소재지 수요 수축까지 3중 증거."),
    ("신용 스트레스가 커지면\n더 깊이 깎는다", BLUE,
     "회사채 스프레드 확대 국면에 저감 심화 —\n1차 차분 상관 -0.68 (p=0.004).\n지가 상승 지역일수록 공매는 빨리 끝난다."),
]
from matplotlib.patches import Rectangle
cw, ch, gap = 0.201, 0.30, 0.018
for i, (title, c, body) in enumerate(cards):
    x = 0.055 + i * (cw + gap)
    card(x, 0.27, cw, ch)
    ax.add_patch(Rectangle((x + 0.015, 0.494), 0.006, 0.05, fc=c, ec="none", zorder=2))
    ax.text(x + 0.033, 0.545, title, color=INK, fontsize=16, weight="bold", va="top")
    ax.text(x + 0.015, 0.44, body, color=INK2, fontsize=12, va="top", linespacing=1.6)

# ── 우측: 유찰 곡선 모티프 ──
axc = fig.add_axes([0.72, 0.60, 0.235, 0.22])
axc.set_facecolor(BG)
rounds = list(range(1, 11))
path = [120 * (0.905 ** (r - 1)) for r in rounds]
axc.plot(rounds, path, color=BLUE, lw=2.5, marker="o", ms=7,
         markerfacecolor=BLUE, markeredgecolor=BG, markeredgewidth=1.5, zorder=3)
axc.axhline(100, color=MUTED, lw=0.8, ls=(0, (4, 4)), alpha=0.6)
axc.scatter([8], [path[7]], color=ORANGE, s=110, zorder=4,
            edgecolor=BG, linewidth=1.5)
axc.annotate("8회차 = 감정가의 62%", (8, path[7]), xytext=(4.6, 48),
             color=ORANGE, fontsize=13, weight="bold")
axc.text(1, 126, "감정가 대비 최저입찰가 경로(회차별, 실측 중앙)", color=INK2, fontsize=13)
axc.set_ylim(40, 140); axc.set_xlim(0.5, 10.5)
for s in axc.spines.values():
    s.set_visible(False)
axc.tick_params(colors=MUTED, labelsize=11, length=0)
axc.set_yticks([60, 80, 100, 120]); axc.set_yticklabels(["60%", "80%", "100%", "120%"])
axc.set_xticks([1, 4, 8]); axc.set_xticklabels(["1회차", "4회차", "8회차"])
axc.grid(axis="y", color="#2A3348", lw=0.6)

# ── 우측 하단: QR + 앱 안내 ──
card(0.72, 0.27, 0.235, 0.30)
qr_ax = fig.add_axes([0.732, 0.295, 0.115, 0.245])
qr_ax.imshow(plt.imread(OUT / "qr_app.png"), interpolation="nearest")
qr_ax.axis("off")
ax.text(0.862, 0.525, "라이브 대시보드", color=INK, fontsize=17, weight="bold", va="top")
for j, feat in enumerate(["지역·유형 필터", "가격 시계열 탐색", "위험 사업장 스크리닝", "사내 데이터 결합 준비"]):
    ax.text(0.862, 0.478 - j * 0.042, "✓", color=AQUA, fontsize=14, weight="bold")
    ax.text(0.878, 0.478 - j * 0.042, feat, color=INK2, fontsize=13.5)
ax.text(0.7325, 0.283, "QR 스캔 또는 링크 접속 — 로그인 불필요", color=MUTED, fontsize=11.5)

# ── 푸터 ──
ax.plot([0.055, 0.945], [0.20, 0.20], color="#2A3348", lw=1)
ax.text(0.055, 0.155, "데이터  전국은행연합회 PF 정보공개 · 온비드 · 국토교통부 실거래가/건축HUB · "
        "한국부동산원 R-ONE · ECOS/KOSIS · 청약홈 · 키스콘 · HUG 등 공개 자료 10여 종",
        color=INK2, fontsize=13.5)
ax.text(0.055, 0.115, "제작  AI 코딩 에이전트 협업 — 생성 → 상호 적대 검수(7차례) → 사람 판정 루프 · "
        "재현 가능한 단일 정본 수치 체계 · 전 과정 익명화",
        color=INK2, fontsize=13.5)
ax.text(0.055, 0.062, URL, color=BLUE, fontsize=14, weight="bold")
ax.text(0.945, 0.062, "PF After-Care Dashboard", color=MUTED, fontsize=13, ha="right")

fig.savefig(OUT / "poster.png", facecolor=BG)
print("poster.png saved", (OUT / "poster.png").stat().st_size // 1024, "KB")

# ═══════════════════ PPTX (편집용, 같은 디자인) ═══════════════════
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def rgb(h):
    return RGBColor.from_string(h.lstrip("#"))

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

def box(x, y, w, h, color=CARD, rounded=True):
    from pptx.enum.shapes import MSO_SHAPE
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = rgb(color)
    shp.line.fill.background()
    return shp

def text(x, y, w, h, s, size, color=INK, bold=False, align=PP_ALIGN.LEFT, sp=1.0):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    lines = s.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = sp
        r = p.add_run(); r.text = ln
        r.font.size = Pt(size); r.font.bold = bold
        r.font.color.rgb = rgb(color); r.font.name = "Malgun Gothic"
    return tb

box(0, 0, 13.333, 7.5, BG, rounded=False)  # 배경
text(0.7, 0.35, 8, 0.35, "공개 데이터 기반 · AI 협업 프로젝트", 13, AQUA, bold=True)
text(0.7, 0.75, 11.9, 0.85, "부실 PF 사업장, 가격은 어떻게 움직이는가", 34, INK, bold=True)
text(0.7, 1.5, 11.9, 0.65,
     "공개 경공매 사업장 798곳의 감정가·최저입찰가 경로를 시계열로 추적하고, 점검이 필요한 사업장을 스크리닝하는 대시보드",
     14, INK2)
stats = [("798", "추적 사업장", BLUE), ("17개월", "가격 경로 관측", BLUE),
         ("129 → 65", "검증으로 확정한 변수", ORANGE), ("7,977건", "통계 검정(FDR)", ORANGE)]
for i, (num, lab, c) in enumerate(stats):
    text(0.7 + i * 2.1, 2.25, 2.0, 0.55, num, 25, c, bold=True)
    text(0.7 + i * 2.1, 2.78, 2.0, 0.3, lab, 11, MUTED)
cards_p = [
    ("가격은 버티고, 최저입찰가만 깎인다", ORANGE,
     "감정평가액은 공개기간 중 83.6% 무변동.\n공매는 감정가의 120%에서 출발해 회차마다\n-9.5%씩 인하 — 8회차엔 62%."),
    ("산업시설은 가격이 아니라 수요의 문제", AQUA,
     "물류·공장은 시장보다 싸게 내놓아도 체류\n최장(중앙 11개월). 감정 괴리 최소·소재지\n수요 수축까지 3중 증거."),
    ("신용 스트레스가 커지면 더 깊이 깎는다", BLUE,
     "회사채 스프레드 확대 국면에 저감 심화 —\n차분 상관 -0.68 (p=0.004). 지가 상승\n지역일수록 공매는 빨리 끝난다."),
]
for i, (title, c, body) in enumerate(cards_p):
    x = 0.7 + i * 2.95
    box(x, 3.35, 2.75, 2.35)
    bar = box(x + 0.18, 3.55, 0.09, 0.5, c)
    text(x + 0.4, 3.55, 2.2, 0.8, title, 13.5, INK, bold=True)
    text(x + 0.18, 4.45, 2.45, 1.2, body, 10.5, INK2, sp=1.2)
box(9.7, 3.35, 3.0, 2.35)
pic = slide.shapes.add_picture(str(OUT / "qr_app.png"), Inches(9.9), Inches(3.55),
                               Inches(1.55), Inches(1.55))
text(11.55, 3.6, 1.1, 0.4, "라이브\n대시보드", 12, INK, bold=True)
text(11.55, 4.35, 1.15, 1.2, "✓ 필터\n✓ 시계열\n✓ 스크리닝", 10, INK2, sp=1.25)
text(9.9, 5.25, 2.7, 0.3, "QR 스캔 — 로그인 불필요", 9, MUTED)
text(0.7, 6.1, 12, 0.3, "데이터  은행연합회 PF 정보공개 · 온비드 · 국토부 실거래가/건축HUB · 부동산원 R-ONE · ECOS/KOSIS · 청약홈 · 키스콘 · HUG 등 공개 자료 10여 종", 10, INK2)
text(0.7, 6.45, 12, 0.3, "제작  AI 코딩 에이전트 협업 — 생성 → 상호 적대 검수(7차례) → 사람 판정 루프 · 전 과정 익명화", 10, INK2)
text(0.7, 6.9, 9, 0.3, URL, 11, BLUE, bold=True)
text(9.0, 6.9, 3.63, 0.3, "PF After-Care Dashboard", 10, MUTED, align=PP_ALIGN.RIGHT)

prs.save(OUT / "poster.pptx")
print("poster.pptx saved")
